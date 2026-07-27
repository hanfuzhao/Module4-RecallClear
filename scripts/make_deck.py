"""Generate the pitch deck as PowerPoint (docs/slides.pptx).

Builds the same eight slides as ``docs/slides.html`` in the same visual
language -- asphalt background, license-plate lockup, hazard stripes,
road-sign colors -- so the .pptx submission and the recorded browser deck are
interchangeable. Every number shown here traces to
``data/outputs/evaluation.json``.

Requires ``python-pptx`` (not part of requirements.txt; it is a documentation
tool, not a pipeline dependency):

    pip install python-pptx
    python -m scripts.make_deck
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from scripts import config

# Palette (matches static/css/styles.css).
ASPHALT = RGBColor(0x1B, 0x1D, 0x21)
NEAR_BLACK = RGBColor(0x11, 0x11, 0x11)
PAINT = RGBColor(0xF7, 0xB5, 0x00)
PAPER = RGBColor(0xFB, 0xF7, 0xEC)
INK_ON_PAPER = RGBColor(0x26, 0x23, 0x19)
TEXT = RGBColor(0xED, 0xEC, 0xE7)
SOFT = RGBColor(0x9A, 0xA0, 0xA8)
FAINT = RGBColor(0x6F, 0x74, 0x7C)
RED = RGBColor(0xD6, 0x40, 0x36)
GREEN = RGBColor(0x2E, 0x9E, 0x63)

CONDENSED = "Arial Narrow"
MONO = "Courier New"
BODY = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(1.1)

OUTPUT_PATH = config.PROJECT_ROOT / "docs" / "slides.pptx"


# --------------------------------------------------------------------------- #
# Building blocks
# --------------------------------------------------------------------------- #


def add_background(slide) -> None:
    """Paint the asphalt background."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ASPHALT


def add_hazard_stripe(slide, top: Emu, height: Emu = Inches(0.12)) -> None:
    """Draw one black-and-yellow hazard stripe across the full width."""
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    base.fill.solid()
    base.fill.fore_color.rgb = PAINT
    base.line.fill.background()
    base.shadow.inherit = False

    step = Inches(0.42)
    x = -step
    while x < SLIDE_W + step:
        block = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, x, top, Emu(int(step * 0.55)), height)
        block.fill.solid()
        block.fill.fore_color.rgb = NEAR_BLACK
        block.line.fill.background()
        block.shadow.inherit = False
        block.adjustments[0] = 0.9
        x += step


def add_text(
    slide,
    left,
    top,
    width,
    height,
    runs,
    font: str = BODY,
    size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.15,
    space_after: int = 6,
):
    """Add a text box. ``runs`` is a list of paragraphs; each paragraph is a
    list of ``(text, overrides)`` run tuples, where overrides may set
    ``bold``, ``color``, ``font``, ``size``."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for p_index, paragraph_runs in enumerate(runs):
        paragraph = frame.paragraphs[0] if p_index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(space_after)
        for text, overrides in paragraph_runs:
            run = paragraph.add_run()
            run.text = text
            run.font.name = overrides.get("font", font)
            run.font.size = Pt(overrides.get("size", size))
            run.font.bold = overrides.get("bold", bold)
            run.font.color.rgb = overrides.get("color", color)
    return box


def plain(text: str, **overrides):
    """One run with optional style overrides."""
    return (text, overrides)


def add_kicker(slide, text: str, badge: str | None = None) -> None:
    """Small yellow section label, optionally with a red innovation badge."""
    runs = [[plain(text.upper(), color=PAINT)]]
    add_text(slide, MARGIN, Inches(0.75), Inches(8.5), Inches(0.4), runs,
             font=CONDENSED, size=16, bold=True)
    if badge:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN + Inches(4.0), Inches(0.70), Inches(4.6), Inches(0.38)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RED
        shape.line.fill.background()
        shape.shadow.inherit = False
        frame = shape.text_frame
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = frame.paragraphs[0].add_run()
        run.text = badge.upper()
        run.font.name = CONDENSED
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_headline(slide, lines: list[list[tuple]], top=Inches(1.15), size: int = 48) -> None:
    """Big condensed uppercase headline."""
    add_text(slide, MARGIN, top, SLIDE_W - MARGIN * 2, Inches(1.9), lines,
             font=CONDENSED, size=size, bold=True, line_spacing=1.02)


def add_foot(slide, left_text: str, right_text: str) -> None:
    """Monospace footer line."""
    add_text(slide, MARGIN, SLIDE_H - Inches(0.62), Inches(6.5), Inches(0.35),
             [[plain(left_text)]], font=MONO, size=11, color=FAINT)
    add_text(slide, SLIDE_W - MARGIN - Inches(6.5), SLIDE_H - Inches(0.62), Inches(6.5), Inches(0.35),
             [[plain(right_text)]], font=MONO, size=11, color=FAINT, align=PP_ALIGN.RIGHT)


def add_plate(slide, top=Inches(0.85)) -> None:
    """The license-plate lockup."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, top, Inches(2.3), Inches(0.62))
    shape.adjustments[0] = 0.25
    shape.rotation = -2
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAINT
    shape.line.color.rgb = RGBColor(0x15, 0x15, 0x0F)
    shape.line.width = Pt(2.5)
    shape.shadow.inherit = False
    frame = shape.text_frame
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = frame.paragraphs[0].add_run()
    run.text = "RECAL·CLR"
    run.font.name = CONDENSED
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x15, 0x15, 0x0F)


def add_table_rows(slide, top, col_lefts, col_widths, rows, row_height=Inches(0.5)) -> None:
    """Lay out a table as aligned text rows (full control over styling).

    ``rows``: list of row definitions; each cell is a (text, style) tuple where
    style may set color/bold/font/size. The first row is the header.
    """
    y = top
    for row_index, row in enumerate(rows):
        for cell_index, (text, style) in enumerate(row):
            is_header = row_index == 0
            add_text(
                slide, col_lefts[cell_index], y, col_widths[cell_index], row_height,
                [[plain(
                    text.upper() if is_header else text,
                    color=style.get("color", SOFT if is_header else TEXT),
                    bold=style.get("bold", is_header),
                    font=style.get("font", CONDENSED if is_header else (MONO if cell_index else BODY)),
                    size=style.get("size", 14 if is_header else 17),
                )]],
                line_spacing=1.0,
            )
        y += row_height


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #


def slide_1_hook(deck) -> None:
    """The real letter, and the measured reading level."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_kicker(slide, "A real letter, mailed to real people")

    paper = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.5), SLIDE_W - MARGIN * 2, Inches(2.5)
    )
    paper.fill.solid()
    paper.fill.fore_color.rgb = PAPER
    paper.line.fill.background()
    paper.shadow.inherit = False
    frame = paper.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.4)
    frame.margin_right = Inches(0.4)
    frame.margin_top = Inches(0.35)
    paragraph = frame.paragraphs[0]
    paragraph.line_spacing = 1.5
    run = paragraph.add_run()
    run.text = (
        '"Underbody heat and noise insulators may loosen and contact the '
        "aluminum driveshaft, which could damage the driveshaft and cause it "
        'to fracture…"'
    )
    run.font.name = MONO
    run.font.size = Pt(20)
    run.font.color.rgb = INK_ON_PAPER

    stamp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, SLIDE_W - MARGIN - Inches(2.3), Inches(1.68), Inches(1.9), Inches(0.42)
    )
    stamp.rotation = 4
    stamp.fill.background()
    stamp.line.color.rgb = RED
    stamp.line.width = Pt(2)
    stamp.shadow.inherit = False
    stamp_run = stamp.text_frame.paragraphs[0].add_run()
    stamp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    stamp_run.text = "SAFETY RECALL"
    stamp_run.font.name = CONDENSED
    stamp_run.font.size = Pt(13)
    stamp_run.font.bold = True
    stamp_run.font.color.rgb = RED

    add_text(slide, MARGIN, Inches(4.35), SLIDE_W - MARGIN * 2, Inches(0.8), [[
        plain("Measured across 11,591 of these letters: average reading level "),
        plain("grade 12.7", color=PAINT, bold=True, font=MONO, size=24),
        plain(" — college."),
    ]], size=19)
    add_foot(slide, "RecallClear", "Duke AIPI 540 · Module 4")


def slide_2_problem(deck) -> None:
    """One in four cars; the letter answers none of the owner's questions."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_kicker(slide, "The problem")
    add_headline(slide, [[plain("1 IN 4 CARS HAS AN")], [plain("OPEN SAFETY RECALL")]])

    add_text(slide, MARGIN, Inches(3.1), Inches(10.5), Inches(3), [
        [plain("•  The repair is "), plain("always free", color=PAINT, bold=True), plain(" — federal law.")],
        [plain("•  Owners still don't get it fixed.")],
        [plain("•  The letter never answers: "),
         plain("how bad? what do I do? what does it cost?", color=PAINT, bold=True)],
    ], size=22, line_spacing=1.4)
    add_foot(slide, "RecallClear", "the letter is the interface problem")


def slide_3_product(deck) -> None:
    """The five-line card and the live app."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_plate(slide)
    add_headline(slide, [[plain("EVERY LETTER → FIVE PLAIN LINES")]], top=Inches(1.7), size=44)

    lines = [
        ("WHAT'S WRONG", " · one sentence"),
        ("WHAT COULD HAPPEN", " · one sentence"),
        ("HOW URGENT", " · a road-sign answer"),
        ("WHAT TO DO", " · the next step + phone number"),
        ("WHAT IT COSTS", " · nothing. always."),
    ]
    paragraphs = [
        [plain(head, color=PAINT, bold=True), plain(tail)] for head, tail in lines
    ]
    add_text(slide, MARGIN, Inches(2.75), Inches(11), Inches(3.2), paragraphs,
             font=MONO, size=20, line_spacing=1.5)
    add_text(slide, MARGIN, Inches(6.15), Inches(11), Inches(0.5),
             [[plain("Live web app — paste the letter, or just the recall number.")]],
             size=17, color=SOFT)
    add_foot(slide, "recallclear-166936551184.us-central1.run.app", "live now")


def slide_4_how(deck) -> None:
    """Method + innovation 1: rule-built supervision from public data."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_kicker(slide, "How", badge="Innovation 1 · self-supervised from public data")
    add_headline(slide, [[plain("LORA ON A 135M MODEL,")], [plain("TRAINED ON A LAPTOP CPU")]], size=42)

    add_text(slide, MARGIN, Inches(3.15), Inches(6.6), Inches(3.2), [
        [plain("•  "), plain("11,591 real notices", color=PAINT, bold=True),
         plain(" — U.S. DOT open data")],
        [plain("•  Zero hand-labeling: targets built by "),
         plain("auditable rules", color=PAINT, bold=True),
         plain(" from NHTSA's fields + official safety flags")],
        [plain("•  LoRA = "), plain("3.5%", color=PAINT, bold=True),
         plain(" of weights → a "), plain("20 MB", color=PAINT, bold=True), plain(" adapter")],
    ], size=18, line_spacing=1.35)

    add_table_rows(
        slide, Inches(3.2),
        [Inches(8.1), Inches(10.6)], [Inches(2.5), Inches(2.2)],
        [
            [("measured", {}), ("", {})],
            [("training time", {}), ("25 min (CPU)", {})],
            [("CPU vs GPU here", {}), ("CPU 4× faster", {"color": PAINT, "bold": True})],
            [("card on CPU", {}), ("~4 s", {})],
        ],
    )
    add_foot(slide, "SmolLM2-135M-Instruct · LoRA r16 · 4.9M trainable params",
             "everything benchmarked, nothing assumed")


def slide_5_results(deck) -> None:
    """The before/after table on held-out brands."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_kicker(slide, "Before / after — 150 notices, 7 brands the model never saw")
    add_text(slide, MARGIN, Inches(1.1), Inches(11), Inches(1.1), [[
        plain("0% → 63% → ", size=56, bold=True),
        plain("100%", color=PAINT, size=56, bold=True),
    ]], font=CONDENSED)

    header = [("", {}), ("stock", {}), ("+2 examples in prompt", {}), ("fine-tuned", {})]
    rows = [
        header,
        [("card format correct", {}), ("0%", {"color": RED, "bold": True}),
         ("63%", {}), ("100%", {"color": PAINT, "bold": True})],
        [("reading grade (letters: 12.7)", {}), ("13.0", {}), ("8.0", {}),
         ("7.3", {"color": PAINT, "bold": True})],
        [("jargon / 100 words", {}), ("2.5", {}), ("2.3", {}),
         ("0.7", {"color": PAINT, "bold": True})],
        [("prompt tokens", {}), ("427", {}), ("1,334", {}),
         ("427", {"color": PAINT, "bold": True})],
    ]
    add_table_rows(
        slide, Inches(2.6),
        [MARGIN, Inches(5.4), Inches(7.2), Inches(10.3)],
        [Inches(4.2), Inches(1.7), Inches(3.0), Inches(2.0)],
        rows, row_height=Inches(0.55),
    )
    add_text(slide, MARGIN, Inches(5.85), Inches(11), Inches(0.5),
             [[plain("The app has three bays — anyone can re-run this race, live.")]],
             size=17, color=SOFT)
    add_foot(slide, "held out by manufacturer, not at random",
             "macro-F1 reported, accuracy is a trap here")


def slide_6_failure(deck) -> None:
    """Innovation 2: the replicated negative result."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_kicker(slide, "The result I'm proudest of", badge="Innovation 2 · a replicated negative")
    add_headline(slide, [
        [plain("WE TRAINED IT TWICE TO SAY")],
        [plain('"'), plain("STOP DRIVING", color=PAINT), plain('." IT WOULDN’T.')],
    ], size=42)

    add_text(slide, MARGIN, Inches(3.25), Inches(11.2), Inches(3.2), [
        [plain("•  Run 1, uniform data: rare-warning recall "),
         plain("0%", color=RED, bold=True, font=MONO)],
        [plain("•  Run 2, rare cases ×8 (16% of the mix): "),
         plain("still 0%", color=RED, bold=True),
         plain(" — even on letters it trained on")],
        [plain("•  Why: the warning is "), plain("3 tokens out of ~130", color=PAINT, bold=True),
         plain(". Token-averaged loss never feels it.")],
        [plain("•  Class imbalance wasn't the problem. "),
         plain("Token imbalance was.", color=PAINT, bold=True)],
    ], size=19, line_spacing=1.4)
    add_foot(slide, "eval loss 0.05 — and still silent", "replicated, not anecdotal")


def slide_7_fix(deck) -> None:
    """Innovation 3: the three-layer safety design."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_kicker(slide, "The fix", badge="Innovation 3 · three-layer safety design")
    add_headline(slide, [[plain("STOP ASKING THE MODEL")]], size=48)

    rows = [
        [("who makes the do-not-drive call", {}), ("recall", {})],
        [("the fine-tuned model (both versions)", {}), ("0%", {"color": RED, "bold": True})],
        [("15 lines of rules reading the letter", {}), ("73%  = all the text contains", {"color": GREEN})],
        [("NHTSA's official flags", {}), ("100%", {"color": GREEN, "bold": True})],
    ]
    add_table_rows(
        slide, Inches(2.35),
        [MARGIN, Inches(7.0)], [Inches(5.6), Inches(5.0)],
        rows, row_height=Inches(0.55),
    )
    add_text(slide, MARGIN, Inches(4.9), Inches(11.2), Inches(1.6), [
        [plain("THE MODEL WRITES THE PROSE.", size=34, bold=True)],
        [plain("THE ALARM IS NEVER ITS CALL.", color=PAINT, size=34, bold=True)],
    ], font=CONDENSED, line_spacing=1.1)
    add_foot(slide, "all three layers ship in the live app", "0% / 73% / 100%")


def slide_8_close(deck) -> None:
    """Links and the closing line."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    add_background(slide)
    add_hazard_stripe(slide, Emu(0))
    add_hazard_stripe(slide, SLIDE_H - Inches(0.12))
    add_plate(slide)
    add_headline(slide, [
        [plain("20 MB. 25 MINUTES. A LAPTOP.")],
        [plain("GRADE 13 → "), plain("GRADE 7", color=PAINT), plain(".")],
    ], top=Inches(1.7), size=44)

    add_text(slide, MARGIN, Inches(3.6), Inches(11.5), Inches(2.0), [
        [plain("Repo   ", color=PAINT, bold=True, font=MONO),
         plain("github.com/hanfuzhao/Module4-RecallClear — 17 reviewed PRs · 85 tests · CI", font=MONO, size=16)],
        [plain("Model  ", color=PAINT, bold=True, font=MONO),
         plain("huggingface.co/HanfuZhao781/recallclear-smollm2-135m-lora", font=MONO, size=16)],
        [plain("App    ", color=PAINT, bold=True, font=MONO),
         plain("recallclear-166936551184.us-central1.run.app", font=MONO, size=16)],
    ], size=16, line_spacing=1.6)

    add_text(slide, MARGIN, Inches(5.7), Inches(11.5), Inches(0.9),
             [[plain("…and the system knows exactly which decision the model is "),
               plain("not allowed to make.", color=PAINT, bold=True)]],
             size=19, color=SOFT)
    add_foot(slide, "thank you", "hanfuzhao · Duke AIPI 540")


def build_deck(output_path: Path = OUTPUT_PATH) -> Path:
    """Assemble all eight slides and save the .pptx."""
    deck = Presentation()
    deck.slide_width = SLIDE_W
    deck.slide_height = SLIDE_H

    for builder in (
        slide_1_hook, slide_2_problem, slide_3_product, slide_4_how,
        slide_5_results, slide_6_failure, slide_7_fix, slide_8_close,
    ):
        builder(deck)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(output_path))
    print(f"Deck written to {output_path} ({len(deck.slides._sldIdLst)} slides)")
    return output_path


if __name__ == "__main__":
    build_deck()
